"""Office file tools — read PDF, DOCX, XLSX, PPTX."""

from __future__ import annotations

import os
from typing import Any

from pyclaw.agents.tools.base import BaseTool
from pyclaw.agents.types import ToolResult


class ReadPdfTool(BaseTool):
    """Extract text from a PDF file."""

    @property
    def name(self) -> str:
        return "read_pdf"

    @property
    def description(self) -> str:
        return "Extract text content from a PDF file. Optionally specify start/end pages."

    @property
    def parameters(self) -> dict[str, Any]:
        return {
            "type": "object",
            "properties": {
                "file_path": {"type": "string", "description": "Path to the PDF file."},
                "start_page": {"type": "integer", "description": "Start page (1-based, optional)."},
                "end_page": {"type": "integer", "description": "End page (1-based, optional)."},
            },
            "required": ["file_path"],
        }

    async def execute(self, tool_call_id: str, arguments: dict[str, Any]) -> ToolResult:
        file_path = arguments.get("file_path", "")
        start = arguments.get("start_page", 1)
        end = arguments.get("end_page")

        if not os.path.isfile(file_path):
            return ToolResult.text(f"File not found: {file_path}", is_error=True)

        try:
            from pypdf import PdfReader
        except ImportError:
            return ToolResult.text(
                "pypdf required. Install with: pip install 'pyclaw[office]'",
                is_error=True,
            )

        try:
            reader = PdfReader(file_path)
            total = len(reader.pages)
            start_idx = max(0, (start or 1) - 1)
            end_idx = min(total, end or total)

            parts: list[str] = [f"PDF: {os.path.basename(file_path)} ({total} pages)\n"]
            for i in range(start_idx, end_idx):
                text = reader.pages[i].extract_text() or ""
                if text.strip():
                    parts.append(f"--- Page {i + 1} ---\n{text.strip()}")

            content = "\n\n".join(parts)
            if len(content) > 100000:
                content = content[:100000] + "\n\n... (truncated)"
            return ToolResult.text(content)
        except Exception as exc:
            return ToolResult.text(f"Error reading PDF: {exc}", is_error=True)


class ReadDocxTool(BaseTool):
    """Extract text from a Word document."""

    @property
    def name(self) -> str:
        return "read_docx"

    @property
    def description(self) -> str:
        return "Extract paragraphs and tables from a Word (.docx) document."

    @property
    def parameters(self) -> dict[str, Any]:
        return {
            "type": "object",
            "properties": {
                "file_path": {"type": "string", "description": "Path to the .docx file."},
            },
            "required": ["file_path"],
        }

    async def execute(self, tool_call_id: str, arguments: dict[str, Any]) -> ToolResult:
        file_path = arguments.get("file_path", "")
        if not os.path.isfile(file_path):
            return ToolResult.text(f"File not found: {file_path}", is_error=True)

        try:
            from docx import Document
        except ImportError:
            return ToolResult.text(
                "python-docx required. Install with: pip install 'pyclaw[office]'",
                is_error=True,
            )

        try:
            doc = Document(file_path)
            parts: list[str] = [f"DOCX: {os.path.basename(file_path)}\n"]

            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    style = para.style.name if para.style else ""
                    if style.startswith("Heading"):
                        parts.append(f"## {text}")
                    else:
                        parts.append(text)

            for i, table in enumerate(doc.tables):
                parts.append(f"\n--- Table {i + 1} ---")
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    parts.append(" | ".join(cells))

            content = "\n".join(parts)
            if len(content) > 100000:
                content = content[:100000] + "\n\n... (truncated)"
            return ToolResult.text(content)
        except Exception as exc:
            return ToolResult.text(f"Error reading DOCX: {exc}", is_error=True)


class ReadXlsxTool(BaseTool):
    """Extract data from an Excel spreadsheet."""

    @property
    def name(self) -> str:
        return "read_xlsx"

    @property
    def description(self) -> str:
        return "Read data from an Excel (.xlsx) spreadsheet. Optionally target a specific sheet."

    @property
    def parameters(self) -> dict[str, Any]:
        return {
            "type": "object",
            "properties": {
                "file_path": {"type": "string", "description": "Path to the .xlsx file."},
                "sheet_name": {"type": "string", "description": "Sheet name to read (optional, reads all if empty)."},
                "max_rows": {"type": "integer", "description": "Maximum rows to read per sheet (default: 500)."},
            },
            "required": ["file_path"],
        }

    async def execute(self, tool_call_id: str, arguments: dict[str, Any]) -> ToolResult:
        file_path = arguments.get("file_path", "")
        sheet_name = arguments.get("sheet_name", "")
        max_rows = arguments.get("max_rows", 500)

        if not os.path.isfile(file_path):
            return ToolResult.text(f"File not found: {file_path}", is_error=True)

        try:
            from openpyxl import load_workbook
        except ImportError:
            return ToolResult.text(
                "openpyxl required. Install with: pip install 'pyclaw[office]'",
                is_error=True,
            )

        try:
            wb = load_workbook(file_path, read_only=True, data_only=True)
            sheets = [sheet_name] if sheet_name else wb.sheetnames
            parts: list[str] = [f"XLSX: {os.path.basename(file_path)} (sheets: {', '.join(wb.sheetnames)})\n"]

            for sn in sheets:
                if sn not in wb.sheetnames:
                    parts.append(f"Sheet '{sn}' not found.")
                    continue
                ws = wb[sn]
                parts.append(f"\n--- Sheet: {sn} ---")
                for row_count, row in enumerate(ws.iter_rows(values_only=True), 1):
                    cells = [str(c) if c is not None else "" for c in row]
                    parts.append(" | ".join(cells))
                    if row_count >= max_rows:
                        parts.append(f"... (truncated at {max_rows} rows)")
                        break

            wb.close()
            content = "\n".join(parts)
            if len(content) > 100000:
                content = content[:100000] + "\n\n... (truncated)"
            return ToolResult.text(content)
        except Exception as exc:
            return ToolResult.text(f"Error reading XLSX: {exc}", is_error=True)


class ReadPptxTool(BaseTool):
    """Extract text from a PowerPoint presentation."""

    @property
    def name(self) -> str:
        return "read_pptx"

    @property
    def description(self) -> str:
        return "Extract slide text and notes from a PowerPoint (.pptx) presentation."

    @property
    def parameters(self) -> dict[str, Any]:
        return {
            "type": "object",
            "properties": {
                "file_path": {"type": "string", "description": "Path to the .pptx file."},
            },
            "required": ["file_path"],
        }

    async def execute(self, tool_call_id: str, arguments: dict[str, Any]) -> ToolResult:
        file_path = arguments.get("file_path", "")
        if not os.path.isfile(file_path):
            return ToolResult.text(f"File not found: {file_path}", is_error=True)

        try:
            from pptx import Presentation
        except ImportError:
            return ToolResult.text(
                "python-pptx required. Install with: pip install 'pyclaw[office]'",
                is_error=True,
            )

        try:
            prs = Presentation(file_path)
            parts: list[str] = [f"PPTX: {os.path.basename(file_path)} ({len(prs.slides)} slides)\n"]

            for i, slide in enumerate(prs.slides, 1):
                texts: list[str] = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                texts.append(text)

                notes_text = ""
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()

                parts.append(f"--- Slide {i} ---")
                if texts:
                    parts.append("\n".join(texts))
                if notes_text:
                    parts.append(f"  [Notes] {notes_text}")

            content = "\n\n".join(parts)
            if len(content) > 100000:
                content = content[:100000] + "\n\n... (truncated)"
            return ToolResult.text(content)
        except Exception as exc:
            return ToolResult.text(f"Error reading PPTX: {exc}", is_error=True)
